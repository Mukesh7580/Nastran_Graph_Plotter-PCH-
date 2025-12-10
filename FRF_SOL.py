
# 1. TOP: Error handling and imports/////////

import traceback
import sys

try:
    import numpy as np
    import pandas as pd
    import matplotlib.pyplot as plt
    import re
    import os
    import tkinter as tk
    from tkinter import Tk, Button, filedialog, messagebox, Label
    from tkinter import ttk, PhotoImage, Label
    from scipy.ndimage import gaussian_filter1d
    from pptx import Presentation
    from pptx.util import Inches
    import threading
except Exception as e:
    print("Error importing modules:", e)
    traceback.print_exc()
    input("Press Enter to exit...")
    sys.exit(1)

# 2. Middle: App class with logic//////

class GraphExporterApp:
    def __init__(self, master):
        self.master = master
        self.df = None
        master.title("Progress")
        master.geometry("800x600")
        master.resizable(False, False)
        master.configure(bg="white")

        logo = PhotoImage(file =r"C:\Users\HP\Downloads\5044527-removebg-preview.png")
        root.iconphoto(False, logo)

        # Set background image
        self.bg_image = PhotoImage(file=r"C:\Users\HP\Downloads\5044527-removebg-preview.png")
        self.bg_label = tk.Label(master, image=self.bg_image,bg="white")

        self.bg_label.place(x=0, y=0, relwidth=1, relheight=1)
        header_frame = tk.Frame(master, bg="black", height=60)
        header_frame.place(x=0, y=0, relwidth=1)

        # Add header text inside the black bar
        header_label = tk.Label(header_frame, text="Graphs",
                                bg="black", fg="white", font=("Segoe UI", 12, "bold"))
        header_label.pack(side="left", padx=20, pady=15)
        #
        # self.upload_button = Button(master, text="Upload .pch File", command=self.upload_file)
        # self.upload_button.pack(pady=10)
        upload_frame = tk.LabelFrame(master, text="Upload File:", bg="#87CEEB")
        upload_frame.place(x=30, y=70, width=740, height=80)

        self.upload_path_entry = ttk.Entry(upload_frame, width=50)
        self.upload_path_entry.grid(row=0, column=0, padx=10, pady=10)

        self.open_button = ttk.Button(upload_frame, text="Open", command=self.browse_file)
        self.open_button.grid(row=0, column=1, padx=5)

        self.upload_button = ttk.Button(upload_frame, text="Upload", command=self.upload_file)
        self.upload_button.grid(row=0, column=2, padx=5)

        self.upload_progress = ttk.Progressbar(master, mode='determinate')
        self.upload_progress.place(x=550, y=97, width=200, height=20)


        # --- Dummy Curve Options & Save ---
        tk.Label(master, text="Curve Width:", bg="#87CEEB").place(x=30, y=160)
        self.curve_width_var = tk.DoubleVar(value=2.0)
        tk.Entry(master, textvariable=self.curve_width_var).place(x=150, y=160)

        tk.Label(master, text="Curve Color:", bg="#87CEED").place(x=30, y=185)
        self.curve_color_var = tk.StringVar(value="red")
        colors = ["red", "blue", "green", "black", "orange"]
        ttk.Combobox(master, textvariable=self.curve_color_var, values=colors, state="readonly").place(x=150, y=185)

        # self.datum_var = tk.BooleanVar()
        # Datum Line Value Input
        tk.Label(master, text="Target value:", bg="#87CEED").place(x=30, y=210)
        self.datum_entry = tk.Entry(master)
        self.datum_entry.insert(0, "0")  # Default value
        self.datum_entry.place(x=150, y=210)


        self.save_button = tk.Button(master, text="Select Save Directory", command=self.save_graphs)
        self.save_button.place(x=30, y=250)

        self.status = tk.Label(master, text="", bg="#E8F0FE")
        self.status.place(x=30, y=600)


    def browse_file(self):
        file_path = filedialog.askopenfilename(title="Select a .pch file", filetypes=[("PCH Files", "*.pch")])
        if file_path:
            self.upload_path_entry.delete(0, tk.END)
            self.upload_path_entry.insert(0, file_path)

    def upload_file(self):
        file_path = self.upload_path_entry.get()
        if not file_path:
            messagebox.showwarning("No File", "No file was selected.")
            return

        self.upload_progress['value'] = 0
        self.status.config(text="Uploading...")
        self.master.after(100, lambda: self.simulate_upload(file_path, step=0))

    def simulate_upload(self, file_path, step):
        if step < 100:
            self.upload_progress['value'] = step
            self.master.after(20, lambda: self.simulate_upload(file_path, step + 5))
        else:
            self.status.config(text="Upload complete!")
            # Use thread for parsing to avoid blocking UI
            threading.Thread(target=self.load_and_parse_file, args=(file_path,), daemon=True).start()

    def load_and_parse_file(self, file_path):
        with open(file_path, "r") as f:
            self.data = f.read()

        # Parse data
        point_blocks = re.split(r'\$POINT ID\s*=\s*(\d+)', self.data)[1:]
        point_ids = []

        for i in range(0, len(point_blocks), 2):
            point_id = int(point_blocks[i])
            block = point_blocks[i + 1]
            freqs = re.findall(r'\s*(\d+\.\d+E[+-]\d+)\s+G', block)
            point_ids.extend([point_id] * len(freqs))

        frequencies = re.findall(r'\s*(\d+\.\d+E[+-]\d+)\s+G', self.data)
        values = re.findall(r"\d+\.\d+E[+-]?\d*\s+G\s+([\d.E+-]+)\s+([\d.E+-]+)\s+([\d.E+-]+)", self.data)
        response = re.findall(r'\$(DISPLACEMENTS|VELOCITY|ACCELERATION)', self.data)
        response_types = []
        response_index = 0

        for i in range(0, len(point_blocks), 2):
            block = point_blocks[i + 1]
            freqs = re.findall(r'\s*\d+\.\d+E[+-]\d+\s+G', block)
            response_types.extend([response[response_index]] * len(freqs))
            response_index += 1

        df = pd.DataFrame({
            "Response": response_types,
            "PointID": point_ids,
            "Frequency": frequencies,
        })

        df[['X', 'Y', 'Z']] = pd.DataFrame(values, columns=["X", "Y", "Z"])
        df["Result"] = pd.to_numeric(df["Z"], errors="coerce")
        df["Frequency"] = pd.to_numeric(df["Frequency"], errors="coerce")
        df = df.dropna(subset=["Frequency", "Result"])

        df["Subcase"] = df.groupby(["Frequency", "PointID", "Response"]).cumcount()
        df_agg = df.groupby(["Frequency", "Response", "PointID"]).agg({"Result": "mean"}).reset_index()
        df_result = df_agg.pivot(index="Frequency", columns=["Response", "PointID"], values="Result")

        self.df = df
        self.status.config(text="File uploaded and processed successfully.")

    def save_graphs(self):
        if self.df is None:
            messagebox.showerror("Error", "Please upload a file first.")
            return

        dir_path = filedialog.askdirectory(title="Select Directory to Save PPT and Images")
        if not dir_path:
            return

        prs = Presentation()
        slide_layout = prs.slide_layouts[6]
        save_path = os.path.join(dir_path, "graph_images")
        os.makedirs(save_path, exist_ok=True)

        point_groups = self.df.groupby("PointID")
        point_ids = list(point_groups.groups.keys())

        for i in range(0, len(point_ids), 6):
            slide = prs.slides.add_slide(slide_layout)
            batch_ids = point_ids[i:i + 6]

            for j, point_id in enumerate(batch_ids):
                group = self.df[self.df["PointID"] == point_id]
                group_clean = group.groupby("Frequency")["Result"].max().reset_index()
                smooth_result = gaussian_filter1d(group_clean["Result"], sigma=2)

                fig, ax = plt.subplots()
                curve_color = self.curve_color_var.get()
                curve_width = self.curve_width_var.get()
                # show_datum = self.datum_var.get()

                ax.plot(group_clean["Frequency"], smooth_result, color=curve_color, linewidth=curve_width)
           
                datum_value_str = self.datum_entry.get()
                try:
                    datum_value = float(datum_value_str)
                except ValueError:
                    datum_value = 0

                ax.axhline(datum_value, color='black', linestyle='--', linewidth=2)

                ax.set_title(f"Point ID: {point_id}")
                ax.set_xlabel("Frequency (Hz)")
                ax.set_ylabel("Z Velocity (mm/s)")
                ax.grid(True)
                fig.tight_layout()

                image_path = os.path.join(save_path, f"plot_{point_id}.jpg")
                fig.savefig(image_path, dpi=150)
                plt.close(fig)

                left = Inches(0.5 + (j % 3) * 3.0)
                top = Inches(0.3 + (j // 3) * 1.6)
                height = Inches(1.5)
                width = Inches(3.0)
                slide.shapes.add_picture(image_path, left, top, width=width, height=height)

        pptx_path = os.path.join(dir_path, "Graph_Export_Presentation.pptx")
        prs.save(pptx_path)
        self.status.config(text=f"Graphs and PPT saved to {pptx_path}")
        messagebox.showinfo("Done", f"Graphs and PPT saved successfully to:\n{pptx_path}")


# 3. Bottom: Run the app


if __name__ == "__main__":
    try:
        root = Tk()
        app = GraphExporterApp(root)
        root.mainloop()
    except Exception as e:
        print("Error during execution:", e)
        traceback.print_exc()
        input("Press Enter to exit...")
